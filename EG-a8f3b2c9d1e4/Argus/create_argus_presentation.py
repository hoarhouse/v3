#!/usr/bin/env python3
"""
Create a stunning, premium Argus Product Vision presentation
with Apple-inspired design and bilingual Hungarian/English content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# Color Palette
BACKGROUND_COLOR = RGBColor(15, 23, 42)  # Deep navy/charcoal
PRIMARY_CYAN = RGBColor(6, 182, 212)  # Cyan/teal
SECONDARY_CYAN = RGBColor(34, 211, 238)  # Lighter cyan
SUCCESS_GREEN = RGBColor(34, 197, 94)  # Green
ALERT_RED = RGBColor(239, 68, 68)  # Red
TEXT_PRIMARY = RGBColor(255, 255, 255)  # White
TEXT_SECONDARY = RGBColor(148, 163, 184)  # Muted gray

def set_slide_background(slide):
    """Set dark background for slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND_COLOR

def add_text_box(slide, left, top, width, height, text, font_size=18, 
                 color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.CENTER):
    """Add a text box with specified formatting"""
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri Light"
    p.alignment = alignment
    
    return textbox

def create_presentation():
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # SLIDE 1: TITLE
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_slide_background(slide1)
    
    # Main headline
    add_text_box(slide1, 1, 1.5, 8, 1, 
                "Argus: Collective Cyber Defense",
                font_size=44, bold=True)
    
    # Hungarian subtitle
    add_text_box(slide1, 1, 2.3, 8, 0.5,
                "Kollektív Kibervédelem",
                font_size=24, color=PRIMARY_CYAN)
    
    # Tagline
    add_text_box(slide1, 1, 3.2, 8, 0.5,
                "Built in Hungary. Protecting the World.",
                font_size=20)
    
    # Hungarian tagline
    add_text_box(slide1, 1, 3.7, 8, 0.4,
                "Magyarországon épült. A világot védi.",
                font_size=16, color=TEXT_SECONDARY)
    
    # Date
    add_text_box(slide1, 7, 4.8, 2, 0.4,
                "December 2025",
                font_size=14, color=TEXT_SECONDARY, alignment=PP_ALIGN.RIGHT)
    
    # SLIDE 2: THE PROBLEM
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide2)
    
    add_text_box(slide2, 0.5, 0.3, 9, 0.6,
                "The Problem",
                font_size=36, bold=True, alignment=PP_ALIGN.LEFT)
    
    # Three key points
    add_text_box(slide2, 0.5, 1.2, 9, 0.5,
                "1. Organizations Fight Alone",
                font_size=22, color=PRIMARY_CYAN, alignment=PP_ALIGN.LEFT, bold=True)
    add_text_box(slide2, 0.8, 1.6, 8, 0.4,
                "Same attacks hit thousands sequentially",
                font_size=18, color=TEXT_SECONDARY, alignment=PP_ALIGN.LEFT)
    
    add_text_box(slide2, 0.5, 2.2, 9, 0.5,
                "2. No Early Warning",
                font_size=22, color=PRIMARY_CYAN, alignment=PP_ALIGN.LEFT, bold=True)
    add_text_box(slide2, 0.8, 2.6, 8, 0.4,
                "If org #3 gets hit, orgs #4-200 have no warning",
                font_size=18, color=TEXT_SECONDARY, alignment=PP_ALIGN.LEFT)
    
    add_text_box(slide2, 0.5, 3.2, 9, 0.5,
                "3. Privacy Prevents Sharing",
                font_size=22, color=PRIMARY_CYAN, alignment=PP_ALIGN.LEFT, bold=True)
    add_text_box(slide2, 0.8, 3.6, 8, 0.4,
                "Centralized platforms require exposing sensitive data",
                font_size=18, color=TEXT_SECONDARY, alignment=PP_ALIGN.LEFT)
    
    # Bottom stats bar
    add_text_box(slide2, 0.5, 4.6, 3, 0.5,
                "€9.2T Annual Cybercrime",
                font_size=16, color=SECONDARY_CYAN, alignment=PP_ALIGN.CENTER, bold=True)
    add_text_box(slide2, 3.5, 4.6, 3, 0.5,
                "€4.88M Average Breach",
                font_size=16, color=SECONDARY_CYAN, alignment=PP_ALIGN.CENTER, bold=True)
    add_text_box(slide2, 6.5, 4.6, 3, 0.5,
                "Every 19 Seconds",
                font_size=16, color=SECONDARY_CYAN, alignment=PP_ALIGN.CENTER, bold=True)
    
    # SLIDE 3: THE BIG IDEA
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide3)
    
    add_text_box(slide3, 1, 0.5, 8, 0.8,
                "Herd Immunity for Cybersecurity",
                font_size=36, bold=True)
    
    add_text_box(slide3, 1, 1.2, 8, 0.5,
                "Nyájimmunitás a kiberbiztonságban",
                font_size=22, color=PRIMARY_CYAN)
    
    add_text_box(slide3, 1, 2.2, 8, 0.6,
                "When 30% of a sector participates, everyone benefits",
                font_size=20)
    
    add_text_box(slide3, 1, 2.9, 8, 0.6,
                "Attacks on one organization protect all others within hours",
                font_size=20)
    
    add_text_box(slide3, 1, 4.2, 8, 0.8,
                "\"Ransomware is a pandemic.\nThis is the vaccination program.\"",
                font_size=24, color=SECONDARY_CYAN, bold=True)
    
    # SLIDE 4: HOW ARGUS WORKS
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide4)
    
    add_text_box(slide4, 0.5, 0.3, 9, 0.6,
                "How Argus Works",
                font_size=36, bold=True, alignment=PP_ALIGN.LEFT)
    
    # Five-step flow
    steps = [
        ("1. DETECT", "Agent identifies threat"),
        ("2. ANONYMIZE", "Strip all identifying data"),
        ("3. SHARE", "Send mathematical fingerprint only"),
        ("4. LEARN", "Network learns pattern"),
        ("5. PROTECT", "Rules pushed to all")
    ]
    
    x_pos = 0.5
    for step, desc in steps:
        add_text_box(slide4, x_pos, 1.3, 1.7, 0.4,
                    step,
                    font_size=14, color=PRIMARY_CYAN, bold=True)
        add_text_box(slide4, x_pos, 1.7, 1.7, 0.6,
                    desc,
                    font_size=12, color=TEXT_SECONDARY)
        x_pos += 1.8
    
    add_text_box(slide4, 1, 3.5, 8, 0.8,
                "Your data stays local.\nOnly mathematical patterns shared.",
                font_size=22, bold=True)
    
    # SLIDE 5: THREE PRODUCTS ONE NETWORK
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide5)
    
    add_text_box(slide5, 1, 0.3, 8, 0.6,
                "Three Products, One Network",
                font_size=36, bold=True)
    
    add_text_box(slide5, 1, 0.9, 8, 0.4,
                "Három termék, egy hálózat",
                font_size=20, color=PRIMARY_CYAN)
    
    # Three columns
    products = [
        ("SENTINEL", "Free Security Agent", "Deploy and protect in 15 minutes"),
        ("IMMUNITY INDEX", "Live Dashboard", "Watch sectors turn from red to green"),
        ("THREAT LAB", "Researcher Platform", "Name threats, earn recognition")
    ]
    
    x_pos = 0.5
    for name, subtitle, desc in products:
        add_text_box(slide5, x_pos, 1.8, 3, 0.5,
                    name,
                    font_size=20, color=SECONDARY_CYAN, bold=True)
        add_text_box(slide5, x_pos, 2.3, 3, 0.4,
                    subtitle,
                    font_size=16, bold=True)
        add_text_box(slide5, x_pos, 2.8, 3, 0.8,
                    desc,
                    font_size=14, color=TEXT_SECONDARY)
        x_pos += 3.2
    
    # SLIDE 6: ARGUS SENTINEL
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide6)
    
    add_text_box(slide6, 0.5, 0.3, 9, 0.6,
                "Argus Sentinel",
                font_size=36, bold=True, alignment=PP_ALIGN.LEFT)
    
    add_text_box(slide6, 0.5, 0.9, 9, 0.4,
                "Free Security Agent | Ingyenes biztonsági ügynök",
                font_size=18, color=PRIMARY_CYAN, alignment=PP_ALIGN.LEFT)
    
    features = [
        "• Lightweight (~50MB), deploys in 15 minutes",
        "• Learns your network baseline with local AI",
        "• Detects anomalies and threat patterns",
        "• Shares only anonymized fingerprints",
        "• Receives detection rules from collective"
    ]
    
    y_pos = 1.7
    for feature in features:
        add_text_box(slide6, 0.8, y_pos, 8, 0.5,
                    feature,
                    font_size=18, alignment=PP_ALIGN.LEFT)
        y_pos += 0.5
    
    add_text_box(slide6, 0.5, 4.5, 9, 0.5,
                "\"Your first line of defense.\"",
                font_size=20, color=SECONDARY_CYAN, bold=True, alignment=PP_ALIGN.LEFT)
    
    # SLIDE 7: IMMUNITY INDEX
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide7)
    
    add_text_box(slide7, 0.5, 0.3, 9, 0.6,
                "Immunity Index",
                font_size=36, bold=True, alignment=PP_ALIGN.LEFT)
    
    add_text_box(slide7, 0.5, 0.9, 9, 0.4,
                "Live Dashboard | Élő nyilvános irányítópult",
                font_size=18, color=PRIMARY_CYAN, alignment=PP_ALIGN.LEFT)
    
    features = [
        "• Real-time map of Europe's cyber defense",
        "• Sectors color-coded: Red → Yellow → Green",
        "• 30% participation = Herd immunity achieved",
        "• Public accountability and competition"
    ]
    
    y_pos = 1.7
    for feature in features:
        add_text_box(slide7, 0.8, y_pos, 8, 0.5,
                    feature,
                    font_size=18, alignment=PP_ALIGN.LEFT)
        y_pos += 0.5
    
    add_text_box(slide7, 0.5, 4.5, 9, 0.5,
                "\"Watch the immune system strengthen in real-time.\"",
                font_size=20, color=SECONDARY_CYAN, bold=True, alignment=PP_ALIGN.LEFT)
    
    # SLIDE 8: THREAT LAB
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide8)
    
    add_text_box(slide8, 0.5, 0.3, 9, 0.6,
                "Threat Lab",
                font_size=36, bold=True, alignment=PP_ALIGN.LEFT)
    
    add_text_box(slide8, 0.5, 0.9, 9, 0.4,
                "Researcher Platform | Kutatói platform",
                font_size=18, color=PRIMARY_CYAN, alignment=PP_ALIGN.LEFT)
    
    features = [
        "• Analyze anonymized attack patterns",
        "• Compete on global leaderboard",
        "• Name discoveries (\"The Martinez Variant\")",
        "• Your rules protect thousands instantly",
        "• Earn recognition and conference invites"
    ]
    
    y_pos = 1.7
    for feature in features:
        add_text_box(slide8, 0.8, y_pos, 8, 0.5,
                    feature,
                    font_size=18, alignment=PP_ALIGN.LEFT)
        y_pos += 0.5
    
    add_text_box(slide8, 0.5, 4.5, 9, 0.5,
                "\"Fame, glory, and better security for everyone.\"",
                font_size=20, color=SECONDARY_CYAN, bold=True, alignment=PP_ALIGN.LEFT)
    
    # SLIDE 9: GLOBAL VISION
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide9)
    
    add_text_box(slide9, 1, 0.3, 8, 0.6,
                "Global Vision",
                font_size=36, bold=True)
    
    add_text_box(slide9, 1, 0.9, 8, 0.4,
                "Globális jövőkép",
                font_size=20, color=PRIMARY_CYAN)
    
    phases = [
        ("EUROPE 2026", "Launch from Budapest, NIS2/DORA tailwinds\n5 countries, 200+ orgs"),
        ("NORTH AMERICA 2027", "MSSP partnerships\nHealthcare + Finance focus"),
        ("ASIA-PACIFIC 2028", "Japan, Singapore, Australia\nData sovereignty demand")
    ]
    
    y_pos = 1.8
    for phase, details in phases:
        add_text_box(slide9, 1, y_pos, 2, 0.4,
                    phase,
                    font_size=18, color=SECONDARY_CYAN, bold=True, alignment=PP_ALIGN.LEFT)
        add_text_box(slide9, 3.2, y_pos, 5, 0.6,
                    details,
                    font_size=16, alignment=PP_ALIGN.LEFT)
        y_pos += 1.0
    
    # SLIDE 10: WHY E-GROUP
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide10)
    
    add_text_box(slide10, 1, 0.3, 8, 0.6,
                "Why E-Group",
                font_size=36, bold=True)
    
    add_text_box(slide10, 1, 0.9, 8, 0.4,
                "Miért az E-Group?",
                font_size=20, color=PRIMARY_CYAN)
    
    assets = [
        "• FedX: Proven federated learning (100K+ records)",
        "• MyD Wallet: First eIDAS 2.0 certified in EU",
        "• IOX: Secure communication infrastructure",
        "• 30 Years: Government & healthcare IT expertise",
        "• Network: Pécs, Semmelweis, EU consortiums"
    ]
    
    y_pos = 1.7
    for asset in assets:
        add_text_box(slide10, 1, y_pos, 8, 0.4,
                    asset,
                    font_size=18, alignment=PP_ALIGN.LEFT)
        y_pos += 0.4
    
    add_text_box(slide10, 1, 4.0, 8, 0.4,
                "Hungarian innovation. European infrastructure. Global ambition.",
                font_size=18, bold=True)
    
    add_text_box(slide10, 1, 4.4, 8, 0.3,
                "Magyar innováció. Európai infrastruktúra. Globális ambíció.",
                font_size=16, color=TEXT_SECONDARY)
    
    # SLIDE 11: PRE-LAUNCH STRATEGY
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide11)
    
    add_text_box(slide11, 0.5, 0.3, 9, 0.6,
                "Pre-Launch Strategy",
                font_size=36, bold=True, alignment=PP_ALIGN.LEFT)
    
    strategies = [
        ("1. Sentinel Free", "Seed 50-100 organizations"),
        ("2. Threat Lab Beta", "Recruit top researchers"),
        ("3. Immunity Index Public", "Create visibility"),
        ("4. First Save Story", "Document real protection")
    ]
    
    y_pos = 1.2
    for title, desc in strategies:
        add_text_box(slide11, 0.8, y_pos, 3, 0.4,
                    title,
                    font_size=18, color=PRIMARY_CYAN, bold=True, alignment=PP_ALIGN.LEFT)
        add_text_box(slide11, 3.8, y_pos, 5, 0.4,
                    desc,
                    font_size=18, alignment=PP_ALIGN.LEFT)
        y_pos += 0.6
    
    add_text_box(slide11, 1, 4.3, 8, 0.5,
                "Success metrics: 100 Orgs | 3 Sectors | 6 Months",
                font_size=20, color=SECONDARY_CYAN, bold=True)
    
    # SLIDE 12: CLOSING
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide12)
    
    add_text_box(slide12, 1, 1.0, 8, 0.8,
                "Köszönöm",
                font_size=48, bold=True)
    
    add_text_box(slide12, 1, 1.8, 8, 0.5,
                "Kérdések? | Questions?",
                font_size=24, color=PRIMARY_CYAN)
    
    topics = [
        "• Technical feasibility",
        "• Resource requirements",
        "• Timeline",
        "• Partnerships"
    ]
    
    y_pos = 2.8
    for topic in topics:
        add_text_box(slide12, 3.5, y_pos, 3, 0.4,
                    topic,
                    font_size=18, alignment=PP_ALIGN.LEFT)
        y_pos += 0.4
    
    add_text_box(slide12, 3.5, 4.8, 3, 0.4,
                "E-Group",
                font_size=16, color=TEXT_SECONDARY)
    
    return prs

if __name__ == "__main__":
    print("Creating Argus Product Vision presentation...")
    presentation = create_presentation()
    
    output_path = "/Users/christopherhoar/Desktop/dcfh/EG-a8f3b2c9d1e4/Argus/Argus_Product_Vision.pptx"
    presentation.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print("✨ Premium presentation created successfully!")