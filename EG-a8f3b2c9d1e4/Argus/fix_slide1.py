#!/usr/bin/env python3
"""
Fix Slide 1 text overlapping issue in Argus presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Color Palette
BACKGROUND_COLOR = RGBColor(15, 23, 42)
PRIMARY_CYAN = RGBColor(6, 182, 212)
TEXT_PRIMARY = RGBColor(255, 255, 255)
TEXT_SECONDARY = RGBColor(148, 163, 184)

def add_text_box(slide, left, top, width, height, text, font_size=18, 
                 color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.CENTER):
    """Add a text box with specified formatting"""
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri Light"
    p.alignment = alignment
    
    return textbox

# Load existing presentation
prs = Presentation('/Users/christopherhoar/Desktop/dcfh/EG-a8f3b2c9d1e4/Argus/Argus_Product_Vision.pptx')

# Get first slide
slide1 = prs.slides[0]

# Clear all existing shapes except background
shapes_to_remove = []
for shape in slide1.shapes:
    if hasattr(shape, 'text'):
        shapes_to_remove.append(shape)

for shape in shapes_to_remove:
    sp = shape._element
    sp.getparent().remove(sp)

# Re-add all text elements with proper spacing
# Adjusted vertical positions to prevent overlapping

# Main headline - moved higher up
add_text_box(slide1, 1, 1.0, 8, 0.6, 
            "Argus: Collective Cyber Defense",
            font_size=44, bold=True)

# Hungarian subtitle - with clear spacing below headline
add_text_box(slide1, 1, 1.8, 8, 0.4,
            "Kollektív Kibervédelem",
            font_size=24, color=PRIMARY_CYAN)

# Tagline - middle of slide with breathing room
add_text_box(slide1, 1, 2.8, 8, 0.4,
            "Built in Hungary. Protecting the World.",
            font_size=20)

# Hungarian tagline - below English tagline
add_text_box(slide1, 1, 3.4, 8, 0.3,
            "Magyarországon épült. A világot védi.",
            font_size=16, color=TEXT_SECONDARY)

# Date - bottom right corner
add_text_box(slide1, 7, 4.9, 2, 0.3,
            "December 2025",
            font_size=14, color=TEXT_SECONDARY, alignment=PP_ALIGN.RIGHT)

# Save the fixed presentation
prs.save('/Users/christopherhoar/Desktop/dcfh/EG-a8f3b2c9d1e4/Argus/Argus_Product_Vision.pptx')
print("✅ Slide 1 fixed - text spacing corrected!")
print("Saved to: /Users/christopherhoar/Desktop/dcfh/EG-a8f3b2c9d1e4/Argus/Argus_Product_Vision.pptx")